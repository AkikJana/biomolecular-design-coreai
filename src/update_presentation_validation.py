"""Bring the deck up to Sections 7.15 to 7.18.

The deck stops at the screening tool. Four sections have landed since, and one
of them -- the check against binding measured in a wet lab -- is the most
defensible result in the dissertation and appears nowhere in the presentation.

Three slides, not four. The readout search and the operator work make one point
between them, which is sharper stated once than twice:

  * which inference setting carried Section 7.13's effect
  * two automated searches, and why one was safe and the other was not
  * the readouts against binding that was actually measured

Layout follows what the deck already enforces: the footer sits at 7.08 in and
add_note is a fixed 0.95 in box, so nothing may start below 6.05 in.

Usage:
    python src/update_presentation_validation.py \
        --in final_dissertation_presentation.pptx \
        --out final_dissertation_presentation.pptx
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "src"))

from update_presentation import (  # noqa: E402
    AMBER, BLUE, GREEN, NAVY, TEAL, add_callout, add_chrome, add_note,
    add_stat_card, blank_slide, move_slide, set_lines,
)


def _body(slide, top, lines, size=15):
    box = slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                   Inches(12.20), Inches(0.40 * len(lines) + 0.3))
    set_lines(box, lines, size=size)
    return box


def _table(slide, rows, top, widths, height=1.5, size=12.5):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.60), Inches(top),
                                 Inches(12.20), Inches(height)).table
    for col, w in zip(tbl.columns, widths):
        col.width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            if not para.runs:
                continue
            para.runs[0].font.size = Pt(size)
            para.runs[0].font.bold = (r == 0)
            para.runs[0].font.color.rgb = NAVY
    return tbl


def slide_settings(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Measured — Which Setting Carried It, and for Which Test", page)
    _body(s, 1.18, [
        "Section 7.13 raised sampling steps, recycling and alignment depth "
        "together. Moving one at a time:",
    ])
    _table(s, [
        ("Arm", "steps / recycle / MSA", "Cohen's d", "share of gain",
         "top-1 rank"),
        ("reduced", "10 / 1 / 32", "0.28", "—", "9 of 22"),
        ("sampling only", "200 / 1 / 32", "1.13", "69%", "13 of 22"),
        ("alignment only", "10 / 1 / full", "0.18", "−8%", "8 of 22"),
        ("recycling only", "10 / 3 / 32", "0.22", "−5%", "7 of 22"),
        ("all three", "200 / 3 / full", "1.52", "100%", "17 of 22"),
    ], 1.95, (2.6, 3.0, 2.0, 2.3, 2.3), height=2.4)
    _body(s, 4.55, [
        "Sampling steps alone recover 69% of the gain for a third of the cost — "
        "33 s a fold against 106 s.",
        "Alignment depth and recycling, moved alone, sit below the reduced "
        "baseline on every readout.",
    ])
    add_callout(s, 5.45,
                "The shares sum to 25–55%, not 100%: the settings are synergistic. "
                "A deeper alignment is worth nothing on an unconverged structure.",
                fill=GREEN)
    return s


def slide_two_searches(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Two Automated Searches — One Was Safe, One Was Not", page)
    _body(s, 1.18, [
        "The same machinery applied to two problems that differ in one respect: "
        "whether a candidate can be checked.",
    ])
    add_stat_card(s, 0.60, 1.85, 5.95, "p = 0.425",
                  "readout search: the gain a permutation null also finds", AMBER)
    add_stat_card(s, 6.85, 1.85, 5.95, "max |diff| = 0",
                  "operator rewrite: verified against its own reference", GREEN)
    items = [
        ("Searching readouts on 22 receptors found a combination beating the "
         "headline by +0.018 AUC",
         "a permutation null gains +0.032 on average, and the held-out panel "
         "reverses the gain to −0.029. Refused."),
        ("Searching operator rewrites found two worth 1.46× and 1.19×",
         "bit-exact, and seeded folds before and after give identical "
         "coordinates and identical ipTM. Accepted."),
    ]
    top = 3.60
    for head, detail in items:
        box = s.shapes.add_textbox(Inches(0.60), Inches(top), Inches(12.20),
                                   Inches(0.85))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = head
        r0.font.size = Pt(14.5)
        r0.font.bold = True
        r0.font.color.rgb = NAVY
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = detail
        r1.font.size = Pt(12.5)
        r1.font.color.rgb = BLUE
        top += 1.05
    add_callout(s, 5.75,
                "Search pays where a candidate can be verified independently of "
                "the data that proposed it.", fill=NAVY)
    return s


def slide_wetlab(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Measured — Against Binding That Was Actually Measured", page)
    _body(s, 1.15, [
        "1,320 de novo designs released by Anthropic, binding measured by two "
        "contract research organisations,",
        "on molecules that postdate every model's training. Contamination is "
        "unavailable as an explanation.",
    ])
    add_stat_card(s, 0.60, 2.15, 3.85, "0.943",
                  "this work, in-training panel", BLUE)
    add_stat_card(s, 4.75, 2.15, 3.85, "0.682–0.803",
                  "this work, held-out panel", TEAL)
    add_stat_card(s, 8.90, 2.15, 3.90, "0.626",
                  "interface pLDDT vs measured binding", AMBER)
    _body(s, 4.05, [
        "Ten independently developed predictors all land between 0.62 and 0.67. "
        "The ceiling belongs to the task.",
    ])
    add_callout(s, 4.75,
                "Interface pLDDT beats the binder's whole-chain confidence by "
                "+0.033 AUC [+0.005, +0.061] — Section 7.7's claim, on measured "
                "binding.", fill=GREEN)
    add_note(s, 5.95,
             "The held-out figure is far closer to the wet-lab figure than the "
             "headline is, and still above it. A benchmark quoted without saying "
             "whether the complexes were in training overstates a screen — this "
             "puts a number on by how much, from outside the dissertation.")
    return s


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src", default="final_dissertation_presentation.pptx")
    ap.add_argument("--out", dest="dst", default="final_dissertation_presentation.pptx")
    args = ap.parse_args()

    prs = Presentation(args.src)
    n_before = len(prs.slides)
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)
    if any("Against Binding That Was Actually Measured" in t for t in titles):
        raise SystemExit("these slides are already in the deck; nothing to do")
    try:
        insert_at = next(i for i, t in enumerate(titles) if t.startswith("Future Plan"))
    except StopIteration:
        raise SystemExit("no 'Future Plan' slide found; deck layout has changed")

    for k, build in enumerate([slide_settings, slide_two_searches, slide_wetlab]):
        build(prs, insert_at + k + 1)
        move_slide(prs, len(prs.slides) - 1, insert_at + k)
        print(f"  inserted at position {insert_at + k + 1}")

    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                    and sh.left and sh.left > Inches(12)):
                set_lines(sh, [str(i)], size=10, color=NAVY)

    prs.save(args.dst)
    print(f"wrote {args.dst}: {len(prs.slides)} slides (+{len(prs.slides)-n_before})")


if __name__ == "__main__":
    main()
