"""Turn the mid-semester deck into the final dissertation deck.

Second pass, after update_presentation.py. It does three things:

  1. re-labels the cover and every slide footer for the final report
  2. inserts a "Results at a Glance" statistics slide
  3. inserts the four measured figures, each directly after the slide that
     states the finding in words

Figures come from make_presentation_figures.py, which draws them from the
results artifacts -- so a chart in the deck cannot drift away from the number in
the report.

Usage:
    python src/make_presentation_figures.py
    python src/finalise_presentation.py
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
ASSETS = REPO_ROOT / "reports" / "assets"

NAVY = RGBColor(0x0B, 0x2B, 0x52)
BLUE = RGBColor(0x1E, 0x6F, 0xB8)
TEAL = RGBColor(0x0E, 0x8F, 0x9A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
WARN_BG = RGBColor(0xFD, 0xF3, 0xE0)
WARN_FG = RGBColor(0x8A, 0x52, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_GREY = RGBColor(0x55, 0x55, 0x55)

OLD_FOOTER = "Boltz-Fast · Mid-Semester Review · Akik Jana (2024AB05287)"
NEW_FOOTER = "Boltz-Fast · Final Dissertation Report · Akik Jana (2024AB05287)"


def set_lines(shape, lines, size=16, color=NAVY, bold=False, align=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        para.space_after = Pt(6)


def add_chrome(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(13.33), Inches(1.00))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.00),
                                  Inches(13.33), Inches(0.06))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLUE; rule.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.12),
                                  Inches(12.20), Inches(0.85))
    set_lines(tb, [title], size=26, color=WHITE, bold=True)
    ft = slide.shapes.add_textbox(Inches(0.55), Inches(7.08),
                                  Inches(9.00), Inches(0.30))
    set_lines(ft, [NEW_FOOTER], size=9, color=FOOTER_GREY)
    pn = slide.shapes.add_textbox(Inches(12.23), Inches(7.08),
                                  Inches(0.70), Inches(0.30))
    set_lines(pn, ["0"], size=10, color=FOOTER_GREY)


def blank_slide(prs):
    layout = next((lo for lo in prs.slide_layouts if lo.name == "Blank"),
                  prs.slide_layouts[-1])
    return prs.slides.add_slide(layout)


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[from_idx])
    lst.insert(to_idx, ids[from_idx])


def figure_slide(prs, title, image, caption):
    """A full-bleed chart with one line of interpretation under it."""
    s = blank_slide(prs)
    add_chrome(s, title)
    path = ASSETS / image
    if not path.exists():
        raise SystemExit(f"missing figure {path}; run make_presentation_figures.py")
    # 7.4x4.0in figures -> scale to 9.6in wide, centred
    s.shapes.add_picture(str(path), Inches(1.85), Inches(1.32), width=Inches(9.6))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60),
                             Inches(6.05), Inches(12.20), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = WARN_BG; box.line.fill.background()
    set_lines(box, [caption], size=12.5, color=WARN_FG)
    return s


def add_callout(slide, top, text, fill=NAVY, size=17, height=1.05):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60),
                                 Inches(top), Inches(12.20), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    set_lines(box, [text], size=size, color=WHITE, bold=True)
    return box


def add_note(slide, top, text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60),
                                 Inches(top), Inches(12.20), Inches(0.95))
    box.fill.solid(); box.fill.fore_color.rgb = WARN_BG; box.line.fill.background()
    set_lines(box, [text], size=12.5, color=WARN_FG)


def find_or_add_body(slide, top):
    return slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                    Inches(12.20), Inches(3.20))


def stat_card(slide, left, top, width, big, caption, fill):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                  Inches(top), Inches(width), Inches(1.55))
    card.fill.solid(); card.fill.fore_color.rgb = fill; card.line.fill.background()
    tf = card.text_frame; tf.word_wrap = True; tf.clear()
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = big
    r0.font.size = Pt(21); r0.font.bold = True; r0.font.color.rgb = WHITE
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = caption
    r1.font.size = Pt(10); r1.font.color.rgb = WHITE


def slide_at_a_glance(prs):
    s = blank_slide(prs)
    add_chrome(s, "Results at a Glance")
    row1 = [("99 / 100%", "tests passing\nunder pytest in CI", GREEN),
            ("87.5%", "MLA KV-cache\nreduction (measured)", BLUE),
            ("1502×", "activation saving\n(trained low-rank only)", BLUE),
            ("8.6×", "interface pLDDT vs ipTM\neffect-to-noise", GREEN)]
    row2 = [("0.378", "OPM held-out error\nat rank 32 (at capacity)", RED),
            ("22 / 132", "receptors / complexes\nPTM-clean panel", TEAL),
            ("0.0628", "ipTM run-to-run SD\n(24 × 4 replicates)", RED),
            ("49%", "of re-runs reproduce\nthe headline p < 0.05", RED)]
    for i, (big, cap, col) in enumerate(row1):
        stat_card(s, 0.60 + i * 3.10, 1.35, 2.90, big, cap, col)
    for i, (big, cap, col) in enumerate(row2):
        stat_card(s, 0.60 + i * 3.10, 3.10, 2.90, big, cap, col)

    tb = s.shapes.add_textbox(Inches(0.60), Inches(4.90), Inches(12.20), Inches(1.10))
    set_lines(tb, [
        "Top row — engineering delivered.  Bottom row — measurement outcomes, "
        "which are what redirected the project:",
        "the low-rank OPM is unreachable on pretrained weights, and ipTM does not "
        "rank binders reproducibly — but interface pLDDT does (slides 19–20).",
    ], size=14)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60),
                             Inches(6.05), Inches(12.20), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = WARN_BG; box.line.fill.background()
    set_lines(box, [
        "⚠  Settings throughout are far below Boltz defaults — 10 sampling steps "
        "vs 200, 1 recycling vs 3, MSA depth 32 vs 8192. The settings confound is "
        "stated, not resolved.",
    ], size=12.5, color=WARN_FG)
    return s


def slide_rescore(prs):
    s = blank_slide(prs)
    add_chrome(s, "Measured — What ipTM Discards Is Recoverable")
    add_callout(s, 1.35, "The negative was about ipTM, not about the structures. "
                         "132 folds were still on disk — re-scoring them cost no "
                         "compute at all.", height=0.95)
    set_lines(find_or_add_body(s, 2.45), [
        "Six interface measures, same complexes, same tests. Only one beats a "
        "cognate's OWN scramble — the test that fixes composition and destroys "
        "only order:",
        "•  ipTM  p = 0.42        •  pDockQ  p = 0.80        •  contacts / density "
        "/ buried SASA  p = 0.05–0.45",
        "•  INTERFACE pLDDT  p < 0.0001    rank 1.91 of 4 (chance 2.50)",
        "",
        "It also survives the reproducibility test that demoted ipTM: effect "
        "+3.30 against run-to-run SD 1.92 = 1.72× noise, where ipTM managed "
        "0.20× — an 8.6× better ratio. Reproduces at p < 0.05 in 100% of re-runs.",
    ], size=14.5)
    add_note(s, 5.95,
             "⚠  pDockQ multiplies interface pLDDT by a contact term — and contacts "
             "run the WRONG way here (scrambles make more: 38.5 vs 32.9), cancelling "
             "the signal. Order-sensitivity is established; receptor-specificity "
             "(p = 0.027) is not yet.")
    return s


def slide_localise(prs):
    """The qualification to 7.6 -- half the signal is foldability."""
    s = blank_slide(prs)
    add_chrome(s, "Measured — Where That Signal Actually Comes From")
    add_callout(s, 1.30, "Interface pLDDT responds to sequence order. Does it "
                         "respond to BINDING, or just to the peptide being "
                         "better folded?", height=0.85)

    set_lines(find_or_add_body(s, 2.30), [
        "Splitting the metric by chain — no new folding:",
        "•  peptide side  +5.25   but whole-chain peptide pLDDT  +4.87   →  mostly "
        "FOLDABILITY, not interface",
        "•  RECEPTOR side  +2.38  (p = 6e-5)  →  the receptor's own residues are "
        "placed better when given the right peptide.",
        "   A disordered peptide cannot do that. This term survives the objection.",
        "",
        "Is the response tied to the binding site?  88 folds, alanine scan:",
        "•  interface→Ala 36.77  vs  surface→Ala 38.62   diff −1.84, p = 0.244 — "
        "right direction, not significant",
        "•  80% power would need 123 receptors; every arm drops 7–12 pLDDT, so both "
        "may sit near a floor.",
    ], size=14)

    add_note(s, 5.95,
             "⚠  This QUALIFIES slide 19. Half the pooled signal is peptide "
             "foldability, which slide 19 did not separate. What stands: the "
             "receptor responds to which peptide it is given. Not established: "
             "that the response is localised to the binding site.")
    return s


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src",
                    default=str(REPO_ROOT / "mid_semester_presentation.pptx"))
    ap.add_argument("--out", dest="dst",
                    default=str(REPO_ROOT / "final_dissertation_presentation.pptx"))
    args = ap.parse_args()

    prs = Presentation(args.src)
    if len(prs.slides) != 16:
        raise SystemExit(f"expected the 16-slide deck from update_presentation.py, "
                         f"got {len(prs.slides)}")

    # -- cover ---------------------------------------------------------------
    cover = prs.slides[0]
    for sh in cover.shapes:
        if sh.has_text_frame and "Mid-Semester Review" in sh.text_frame.text:
            set_lines(sh, ["BITS ZG628T Dissertation — Final Report"],
                      size=16, color=RGBColor(0xCF, 0xE0, 0xF0))
            for extra, sz, bold, col in (
                    ("Akik Jana  (2024AB05287)", 15, True, WHITE),
                    ("Supervisor: Dr. Arnab Bandyopadhyay · Dr. Reddy's Laboratories",
                     12.5, False, RGBColor(0xCF, 0xE0, 0xF0)),
                    ("M.Tech. (AI & ML) · BITS Pilani · August 2026",
                     12.5, False, RGBColor(0xCF, 0xE0, 0xF0))):
                para = sh.text_frame.add_paragraph()
                run = para.add_run(); run.text = extra
                run.font.size = Pt(sz); run.font.bold = bold
                run.font.color.rgb = col
                para.space_after = Pt(6)

    # -- footers -------------------------------------------------------------
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and OLD_FOOTER in sh.text_frame.text:
                set_lines(sh, [NEW_FOOTER], size=9, color=FOOTER_GREY)

    # -- reconcile slides written before the measurement phase ---------------
    #
    # Slides 2, 9, 10 and 11 predate Sections 7.1-7.5 and assert things the
    # later slides contradict: an untested screening objective, an affinity
    # surrogate described as "not yet trained" when it was trained and failed,
    # and a "next step" that measurement has since displaced. Cross-references
    # are also stale -- they were written when the deck was 16 slides, and
    # inserting five shifted every target.
    text_fixes = {
        "(see slide 11)": "(see slide 12)",
        "(slides 12–13)": "(slides 14–18)",
        "(slides 11–13)": "(slides 11–18)",
        "Accelerating the full Boltz-2 model is the next objective.":
            "Track A was then tested end-to-end: the ipTM reference does not "
            "rank binders, so the screening objective is not reachable through "
            "it (slides 14–18).",
        "Affinity surrogate head":
            "Affinity surrogate head (since distilled — negative, slide 14)",
        "Therefore the high-value next step is sampler distillation on the "
        "full model — not more weight tricks.":
            "Sampler distillation remains the route to minutes → seconds. "
            "Measurement has since put replicate-averaged folding on GPU ahead "
            "of it in priority (slide 18).",
    }
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    for a, b in text_fixes.items():
                        if a in run.text:
                            run.text = run.text.replace(a, b)

    # -- new slides, appended then moved into place --------------------------
    slide_at_a_glance(prs)                                              # -> 11
    figure_slide(prs, "Low-Rank OPM — Rank vs Fidelity", "fig_opm_rank.png",
                 "Fitted on held-out points from the 33-fold corpus. Reaching 10% "
                 "error needs rank ≈1,414 (layer_0) and ≈9,750 (layer_1) — both "
                 "beyond the width stock actually materialises.")        # -> 13
    figure_slide(prs, "ipTM by Class — the Scramble Control",
                 "fig_iptm_classes.png",
                 "Scrambles sit level with cognates and above decoys. Composition "
                 "is shared with the scramble; order is not — so the separation "
                 "from decoys is compositional.")                        # -> 15
    figure_slide(prs, "Rank Stability Across Identical Re-Runs",
                 "fig_rank_stability.png",
                 "Four receptors, four identical re-runs each. Every one changes "
                 "rank; 6YOO moves from best to worst. Per-receptor rankings are "
                 "not a property of the receptor.")                      # -> 17
    figure_slide(prs, "Would the Headline Survive a Re-Run?",
                 "fig_pvalue_repro.png",
                 "Parametric bootstrap over the measured noise, 4,000 replications. "
                 "The effect direction is robust; the significance verdict is a "
                 "coin flip.")                                           # -> 18
    slide_rescore(prs)                                                   # -> 19
    figure_slide(prs, "ipTM vs Interface pLDDT — the Scramble Control",
                 "fig_metric_comparison.png",
                 "Same structures, two readouts. ipTM cannot separate a cognate "
                 "from its own scramble; interface pLDDT can. The information is "
                 "in the prediction — ipTM discards it.")                # -> 20
    slide_localise(prs)                                                  # -> 21

    n = len(prs.slides)
    # appended: glance, opm, iptm, rank, pval, rescore, rescore-fig, localise
    for src_idx, dst_idx in ((n - 8, 11),    # glance   -> after Key Insight
                             (n - 7, 13),    # opm fig  -> after OPM text
                             (n - 6, 15),    # iptm fig -> after ipTM text
                             (n - 5, 17),    # rank fig -> after repro text
                             (n - 4, 18),    # pval fig
                             (n - 3, 19),    # rescore text
                             (n - 2, 20),    # rescore figure
                             (n - 1, 21)):   # signal localisation
        move_slide(prs, src_idx, dst_idx)

    # -- renumber ------------------------------------------------------------
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.left and sh.left > Inches(12)
                    and sh.top and sh.top > Inches(7)):
                set_lines(sh, [str(i - 1)], size=10, color=FOOTER_GREY)

    prs.save(args.dst)
    print(f"wrote {args.dst}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
