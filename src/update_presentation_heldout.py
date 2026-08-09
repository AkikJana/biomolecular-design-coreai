"""Bring the deck up to the report's Sections 7.10 and 7.11.

The deck stops at Section 7.8 and still presents "interface pLDDT ranks binders
where ipTM does not" as a headline. Two things happened since:

  * a 22-receptor held-out panel, folded three times, shows both readouts losing
    about half their effect on complexes the model was not trained on -- and
    retracts two earlier claims that the loss was specific to one metric
  * auditing the coordinates shows that at ten sampling steps only 14% of
    backbone bonds are physically plausible, against 96% for a distilled model,
    which explains four separate readout failures at once

Four slides are inserted before Future Plan, and the two stale headline claims
are corrected in place.

Styling helpers are imported from update_presentation.py so the deck keeps one
visual language; shape lookups are by text so this fails loudly rather than
silently rewriting the wrong box.

Usage:
    python src/update_presentation_heldout.py \
        --in final_dissertation_presentation.pptx \
        --out final_dissertation_presentation.pptx
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "src"))

from update_presentation import (  # noqa: E402
    AMBER, BLUE, GREEN, NAVY, TEAL, add_callout, add_chrome, add_note,
    add_stat_card, blank_slide, move_slide, set_lines,
)


def _body(slide, top, lines, size=15):
    box = slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                   Inches(12.20), Inches(0.40 * len(lines) + 0.3))
    set_lines(box, lines, size=size)
    return box


def slide_heldout(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Measured — Held-Out Complexes Cost Half the Effect", page)
    _body(s, 1.30, [
        "22 receptors released after Boltz-1's 2021-09-30 training cutoff, screened identically,",
        "decoys drawn from within the held-out set. Folded three times on DeCAF.",
    ])
    add_stat_card(s, 0.60, 2.20, 3.85, "+12.03 → +4.99",
                  "interface pLDDT, in-training → held out", BLUE)
    add_stat_card(s, 4.75, 2.20, 3.85, "+0.265 → +0.137",
                  "ipTM, in-training → held out", TEAL)
    add_stat_card(s, 8.90, 2.20, 3.90, "≈ 2×",
                  "how optimistic an unsplit benchmark is", AMBER)
    _body(s, 4.10, [
        "Both readouts retain roughly half — interface pLDDT 41%, ipTM 52%.",
        "Neither can be shown to degrade more than the other on this panel.",
    ])
    add_callout(s, 5.05,
                "A screening figure quoted without saying whether the complexes "
                "were in training is about twice as good as it should be.")
    add_note(s, 6.20,
             "Not homology-decontaminated, and cannot be: 0 of 22 held-out receptors clear "
             "30% identity to the pre-cutoff PDB (median 1.000), and neither does a random "
             "sample of the whole candidate pool. What the split isolates is complex-level "
             "novelty with receptor familiarity held constant.")
    return s


def slide_draws(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Three Draws — and Two Retracted Claims", page)
    _body(s, 1.30, [
        "Folds are unseeded, so re-running the identical panel gives an independent draw.",
    ])
    rows = [
        ("Cognate rank among its own decoys (chance 2.50)", "draw 1", "draw 2", "draw 3", "mean"),
        ("ipTM", "1.64", "1.50", "1.86", "1.73  (p = 0.007)"),
        ("interface pLDDT", "2.09", "2.05", "1.73", "1.91  (p = 0.020)"),
    ]
    tbl = s.shapes.add_table(len(rows), 5, Inches(0.60), Inches(1.95),
                             Inches(12.20), Inches(1.5)).table
    for c, w in zip(tbl.columns, (4.4, 1.7, 1.7, 1.7, 2.7)):
        c.width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(13)
            para.runs[0].font.bold = (r == 0)
            para.runs[0].font.color.rgb = NAVY
    _body(s, 3.75, [
        "Draw 1 alone said interface pLDDT collapses held out.",
        "Draws 1 and 2 said ipTM specifically survives and interface pLDDT specifically fails.",
        "Draw 3 reverses the ordering. Averaged, the two differ by less than one draw's spread.",
    ])
    add_callout(s, 5.30,
                "Both claims withdrawn. Section 7.5's own finding, turned on this "
                "project's headline — twice.", fill=AMBER)
    add_note(s, 6.45,
             "The methodological lesson is sharper than the scientific one: a single unseeded "
             "fold is not a measurement, and on this evidence neither are two.")
    return s


def slide_geometry(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Measured — At Ten Steps the Coordinates Are Not a Protein", page)
    _body(s, 1.30, [
        "A peptide bond fixes consecutive alpha-carbons at 3.80 Å. Measuring every "
        "consecutive CA–CA distance:",
    ])
    add_stat_card(s, 0.60, 2.05, 5.95, "14.0%",
                  "Boltz-2 at 10 steps — physically plausible backbone bonds", AMBER)
    add_stat_card(s, 6.85, 2.05, 5.95, "96.2%",
                  "DeCAF, distilled FOR 10 steps — same measurement", GREEN)
    _body(s, 3.95, [
        "Boltz-2: median CA–CA 5.48 Å, 56% of bonds beyond 5 Å, 99% of chains majority-broken.",
        "DeCAF:   median 3.74 Å, zero broken bonds.  Residue numbering is sequential —",
        "         this is not a file-ordering artefact. The backbone is genuinely not connected.",
    ])
    add_callout(s, 5.45,
                "Section 7.8 measured that distillation buys a 5–6× larger effect. "
                "This is what it bought: geometry.", fill=GREEN)
    return s


def slide_one_cause(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "One Cause Behind Four Separate Failures", page)
    _body(s, 1.25, [
        "Geometry-dependent readouts need a converged sampler. Four results now share one cause:",
    ], size=15)
    items = [
        ("pDockQ's contact term ran backwards",
         "scrambles made more contacts than cognates — reverses on converged structures "
         "(cognate 61.4, scramble 51.8)"),
        ("pDockQ2 repairs it",
         "swapping the contact term for a PAE term: scramble control p = 0.797 → 0.00026"),
        ("Minimum PAE failed here but is published as best-in-class",
         "fails on Boltz-2 (p = 0.611), among the best on DeCAF (rank 1.55) — PAE predicts "
         "geometry, and was published at full sampling"),
        ("Physics rescoring returned femtomolar affinities",
         "PRODIGY on Boltz-2 structures: garbage in, not method failure"),
    ]
    top = 1.95
    for head, detail in items:
        box = s.shapes.add_textbox(Inches(0.60), Inches(top), Inches(12.20), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = head
        r0.font.size = Pt(15); r0.font.bold = True; r0.font.color.rgb = NAVY
        p1 = tf.add_paragraph()
        r1 = p1.add_run(); r1.text = detail
        r1.font.size = Pt(12.5); r1.font.color.rgb = BLUE
        top += 1.12
    add_note(s, 6.55,
             "Physics rescoring on converged structures still does not help: PRODIGY ΔG reaches "
             "AUC 0.563 against 0.856 for interface pLDDT, and combining them costs 0.033. "
             "A sequence model trained on 96M interactions (MINT) reaches 0.614 zero-shot — "
             "structure wins here, contrary to recent affinity-regression results.")
    return s


def correct_stale_claims(prs):
    """Fix the two headline claims the held-out panel qualifies."""
    fixed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if "Interface pLDDT Ranks Binders Where ipTM Does Not" in txt:
                set_lines(shape, ["Measured — What ipTM Discards Is Recoverable"],
                          size=26, color=NAVY, bold=True)
                fixed += 1
            elif "What ipTM Discards Is Recoverable" in txt and len(txt) < 60:
                fixed += 0  # already the corrected wording
    return fixed


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src", default="final_dissertation_presentation.pptx")
    ap.add_argument("--out", dest="dst", default="final_dissertation_presentation.pptx")
    args = ap.parse_args()

    prs = Presentation(args.src)
    n_before = len(prs.slides)
    print(f"opened {args.src}: {n_before} slides")

    # Future Plan and Summary are the last two; new slides go before them.
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)
    try:
        insert_at = next(i for i, t in enumerate(titles) if t.startswith("Future Plan"))
    except StopIteration:
        raise SystemExit("no 'Future Plan' slide found; deck layout has changed")

    builders = [slide_heldout, slide_draws, slide_geometry, slide_one_cause]
    for k, build in enumerate(builders):
        build(prs, insert_at + k + 1)
        move_slide(prs, len(prs.slides) - 1, insert_at + k)
        print(f"  inserted at position {insert_at + k + 1}")

    fixed = correct_stale_claims(prs)
    print(f"  corrected {fixed} stale headline claim(s)")

    # renumber the trailing slides that shifted
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                    and sh.left and sh.left > Inches(12)):
                set_lines(sh, [str(i)], size=10, color=NAVY)

    prs.save(args.dst)
    print(f"wrote {args.dst}: {len(prs.slides)} slides "
          f"(+{len(prs.slides) - n_before})")


if __name__ == "__main__":
    main()
