"""Bring the mid-semester deck up to date with the measurement phase.

The deck was built before any end-to-end measurement existed: it reports 53
tests, presents the low-rank OPM saving without the reconstruction error, and
lists "distil the affinity surrogate and quantify ranking agreement" as future
work -- which has since been done, and failed. This script applies the updates
and inserts the three measured results.

It edits the existing deck rather than regenerating it, so the hand-built layout
and styling survive. Shape lookups are by text content rather than index, so the
script fails loudly if the deck changes underneath it instead of silently
rewriting the wrong box.

Usage:
    python src/update_presentation.py --in <deck.pptx> --out <deck.pptx>
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]

NAVY = RGBColor(0x0B, 0x2B, 0x52)
BLUE = RGBColor(0x1E, 0x6F, 0xB8)
TEAL = RGBColor(0x0E, 0x8F, 0x9A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xC6, 0x28, 0x28)
WARN_BG = RGBColor(0xFD, 0xF3, 0xE0)
WARN_FG = RGBColor(0x8A, 0x52, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER = RGBColor(0x55, 0x55, 0x55)

FOOTER_TEXT = "Boltz-Fast · Mid-Semester Review · Akik Jana (2024AB05287)"


def find_shape(slide, needle):
    """The one shape whose text contains `needle`, or raise."""
    hits = [s for s in slide.shapes
            if s.has_text_frame and needle in s.text_frame.text]
    if len(hits) != 1:
        raise SystemExit(f"expected 1 shape containing {needle!r}, found {len(hits)}")
    return hits[0]


def set_lines(shape, lines, size=16, color=NAVY, bold=False):
    """Replace a text frame's paragraphs, keeping the deck's type styling."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        para.space_after = Pt(6)


def add_chrome(slide, title, page):
    """Header bar, accent rule, title, footer, page number."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(13.33), Inches(1.00))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.00),
                                  Inches(13.33), Inches(0.06))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLUE; rule.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.12),
                                  Inches(12.20), Inches(0.85))
    set_lines(tb, [title], size=26, color=WHITE, bold=True)

    ft = slide.shapes.add_textbox(Inches(0.55), Inches(7.08),
                                  Inches(9.00), Inches(0.30))
    set_lines(ft, [FOOTER_TEXT], size=9, color=FOOTER)

    pn = slide.shapes.add_textbox(Inches(12.23), Inches(7.08),
                                  Inches(0.70), Inches(0.30))
    set_lines(pn, [str(page)], size=10, color=FOOTER)


def add_stat_card(slide, left, top, width, big, caption, fill):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                  Inches(top), Inches(width), Inches(1.60))
    card.fill.solid(); card.fill.fore_color.rgb = fill; card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = big
    r0.font.size = Pt(23); r0.font.bold = True; r0.font.color.rgb = WHITE
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = caption
    r1.font.size = Pt(10.5); r1.font.color.rgb = WHITE


def add_callout(slide, top, text, fill=NAVY, size=17, height=1.05):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60),
                                 Inches(top), Inches(12.20), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    set_lines(box, [text], size=size, color=WHITE, bold=True)
    return box


def add_note(slide, top, text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60),
                                 Inches(top), Inches(12.20), Inches(0.95))
    box.fill.solid(); box.fill.fore_color.rgb = WARN_BG; box.line.fill.background()
    set_lines(box, [text], size=12.5, color=WARN_FG)


def blank_slide(prs):
    layout = next((lo for lo in prs.slide_layouts if lo.name == "Blank"),
                  prs.slide_layouts[-1])
    return prs.slides.add_slide(layout)


def move_slide(prs, from_idx, to_idx):
    """python-pptx only appends; reorder via the slide id list."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[from_idx])
    sldIdLst.insert(to_idx, ids[from_idx])


# ---------------------------------------------------------------- new slides

def slide_opm(prs):
    s = blank_slide(prs)
    add_chrome(s, "Measured — Low-Rank OPM on Pretrained Weights", 11)
    add_callout(s, 1.42, "The ~97% activation saving is not reachable on "
                         "released checkpoints. The reason is capacity, not "
                         "optimisation or data.")
    set_lines(find_or_add_body(s, 2.70), [
        "Three experiments, each stricter — held-out error at rank 32:",
        "•  CP projection of weights (random inputs) ....................  0.77",
        "•  Fit on one target's activations (one protein) ...............  0.43",
        "•  Corpus distillation, 33 folds (at capacity) .................  0.378",
        "",
        "Train and held-out coincide (0.375 / 0.378) — the signature of a model "
        "at capacity, so more folding data cannot lower the floor.",
    ], size=15)
    add_note(s, 5.75,
             "⚠  Scaling law 1.32·rank^−0.356 puts 10% error at rank ≈1414, "
             "against c_hidden² = 1024 — the low-rank form costs MORE activation "
             "memory than stock before it is accurate enough to substitute.")
    return s


def slide_iptm(prs):
    s = blank_slide(prs)
    add_chrome(s, "Measured — Does ipTM Rank Binders?", 12)
    add_stat_card(s, 0.60, 1.42, 2.95, "22 / 132", "receptors / complexes\n"
                                                   "PTM-clean panel", BLUE)
    add_stat_card(s, 3.75, 1.42, 2.95, "0.5015", "cognate\nmean ipTM", TEAL)
    add_stat_card(s, 6.90, 1.42, 2.95, "0.4888", "SCRAMBLED\nmean ipTM", AMBER)
    add_stat_card(s, 10.05, 1.42, 2.75, "0.4317", "decoy\nmean ipTM", GREEN)
    set_lines(find_or_add_body(s, 3.25), [
        "A scramble keeps composition and length exactly and destroys only "
        "order — and order is what makes a binder a binder.",
        "•  Scrambles outscore genuine binders of other receptors:  AUC 0.632, "
        "p = 0.0096  (110 complexes)",
        "•  Cognate − own scramble:  +0.0128, 95% CI [−0.019, +0.044] — a bound, "
        "not a zero",
        "•  Not a length artefact: the advantage survives length adjustment "
        "(intercept +0.075, p = 0.0001)",
    ], size=15)
    add_note(s, 5.90,
             "⚠  ipTM responds largely to peptide COMPOSITION. Any sequence-order "
             "effect is at most 63% of the composition effect. A ridge regression "
             "on composition independently beat the distilled surrogate "
             "(+0.103 vs −0.034).")
    return s


def slide_repro(prs):
    s = blank_slide(prs)
    add_chrome(s, "Measured — A Single Fold Does Not Reproduce Itself", 13)
    add_callout(s, 1.42, "Boltz's --seed defaults to None. Every benchmark fold "
                         "was a single unseeded draw from a distribution whose "
                         "width had never been measured.", height=1.05)
    set_lines(find_or_add_body(s, 2.70), [
        "24 complexes × 4 identical re-runs (96 folds), same settings:",
        "•  Pooled within-complex ipTM SD  0.0628   (median range 0.127)",
        "•  Cognate − decoy  +0.0698 = 1.11 × SD    comparable to noise",
        "•  Per-receptor rank flips for 4 of 4 receptors — 6YOO spans "
        "#1 to #4 on identical input",
        "•  Aggregate survives: mean rank 2.03, 95% range [1.77, 2.27], always "
        "below chance …",
        "•  … but p < 0.05 in only 49% of simulated re-runs (median p = 0.054)",
    ], size=15)
    add_note(s, 5.95,
             "⚠  Stabilising per-receptor ranks needs ~9–16 replicate folds per "
             "complex (1,200–2,100 folds) — GPU-scale. Ranking binders on a single "
             "AlphaFold/Boltz run is common practice, and does not reproduce.")
    return s


def find_or_add_body(slide, top):
    return slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                    Inches(12.20), Inches(2.90))


# ---------------------------------------------------------------- edits

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src",
                    default=str(REPO_ROOT / "mid_semester_presentation.pptx"))
    ap.add_argument("--out", dest="dst",
                    default=str(REPO_ROOT / "mid_semester_presentation.pptx"))
    args = ap.parse_args()

    prs = Presentation(args.src)
    if len(prs.slides) != 13:
        raise SystemExit(f"expected the 13-slide deck, got {len(prs.slides)}")

    # -- slide 8: benchmarks -------------------------------------------------
    s8 = prs.slides[7]
    card = find_shape(s8, "53 / 100%")
    card.text_frame.paragraphs[0].runs[0].text = "99 / 100%"
    set_lines(find_shape(s8, "Low-rank memory savings are large"), [
        "•  Low-rank memory savings are large and real, and scale with sequence "
        "length — but only for a model trained around the layer (see slide 11).",
        "•  Speculative sampler reduces heavy target evaluations at low tolerance "
        "in component tests.",
    ], size=15)
    set_lines(find_shape(s8, "⚠  Caveats:"), [
        "⚠  Caveats:  (1) low-rank Relative MSE is 84–93% at these ranks — the "
        "layer must be TRAINED to recover full-rank behaviour, and cannot be "
        "projected onto pretrained weights.  (2) Latencies are the surrogate, not "
        "full Boltz-2.  (3) Settings are far below Boltz defaults: 10 sampling "
        "steps vs 200, MSA depth 32 vs 8192.",
    ], size=12, color=WARN_FG)

    # -- slide 9: verification ----------------------------------------------
    s9 = prs.slides[8]
    set_lines(find_shape(s9, "4-tier suite"), [
        "•  99-test suite under pytest in CI, with a ruff lint gate.",
        "•  100% local pass rate; deterministic regression test per component.",
        "•  Equivalence tests where applicable (Fold-CP, OPM S-contraction match "
        "dense to ~1e-7 / 0.0).",
        "–  Scope: component- and synthetic-level; full-model biological "
        "validation is covered by the PDB benchmark (slides 12–13).",
    ], size=15)
    set_lines(find_shape(s9, "Test Tiers"), [
        "Two defects found",
        "•  CI ran `unittest discover`",
        "   against a pytest suite —",
        "   it collected ZERO tests",
        "   and always reported green.",
        "•  Audit found unfailable",
        "   assertions and RMSD",
        "   thresholds satisfied by",
        "   random noise.",
    ], size=12)

    # -- new measured-results slides, after "Key Insight" (index 10) ---------
    for builder in (slide_opm, slide_iptm, slide_repro):
        builder(prs)
    n = len(prs.slides)
    for offset, src_idx in enumerate(range(n - 3, n)):
        move_slide(prs, src_idx, 11 + offset)

    # -- future plan ---------------------------------------------------------
    s_future = prs.slides[14]
    repl = {
        "Profile a full Boltz-2 run": "DONE — profiled; and measured what the "
        "optimisations actually deliver on pretrained weights (slides 11–13).",
        "Few-step diffusion-sampler distillation": "Few-step diffusion-sampler "
        "distillation (consistency/progressive) — keep the full trunk; target "
        "minutes → seconds locally.",
        "Distil the affinity surrogate": "DONE, NEGATIVE — the surrogate was "
        "distilled and evaluated: ipTM does not rank binders, so no surrogate "
        "distilled from it can either.",
        "Use open Boltz-2 weights only": "Use open Boltz-2 weights only — "
        "Boltz-2.1 is closed / API-only.",
        "Production GPU scaling": "Production GPU scaling — now a PREREQUISITE, "
        "not a throughput nicety: replicate-averaged folds (9–16 per complex) "
        "are needed before any per-receptor claim holds.",
    }
    for needle, text in repl.items():
        set_lines(find_shape(s_future, needle), [text], size=13)

    # -- summary -------------------------------------------------------------
    s_sum = prs.slides[15]
    set_lines(find_shape(s_sum, "Implemented & integrated"), [
        "•  Implemented and integrated memory, sampling and parallelism "
        "optimisations into the Boltz pipeline; device-agnostic on Apple Silicon.",
        "•  Moved from component microbenchmarks to end-to-end measurement "
        "against pretrained checkpoints and experimentally determined structures.",
        "•  Two questions closed with negative answers, both with mechanisms: "
        "the low-rank OPM is unreachable on released weights, and ipTM tracks "
        "peptide composition rather than binding.",
        "•  A third finding generalises beyond this project: a single unseeded "
        "fold does not reproduce its own ranking.",
        "•  Next: GPU-scale replicate folding, then few-step sampler distillation.",
    ], size=14)

    # -- renumber page numbers ----------------------------------------------
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.left and sh.left > Inches(12)
                    and sh.top and sh.top > Inches(7)):
                set_lines(sh, [str(i - 1)], size=10, color=FOOTER)

    prs.save(args.dst)
    print(f"wrote {args.dst}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
