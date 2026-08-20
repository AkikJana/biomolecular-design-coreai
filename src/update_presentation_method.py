"""Add the methodological slide: what pre-specification refused.

Section 8.1 now names the practice that produced the later half of Section 7,
and the deck has no equivalent. It shows results but not the discipline that
decided which results survived -- which for a viva is the part that transfers.

The slide leads with what the practice refused rather than with the practice
itself, because four refusals are evidence and a description is not. It goes
immediately before Future Plan so it reads as the closing methodological point.

Usage:
    python src/update_presentation_method.py
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "src"))

from update_presentation import (  # noqa: E402
    GREEN, NAVY, add_callout, add_chrome, add_note, blank_slide,
    move_slide, set_lines,
)


def _body(slide, top, lines, size=15):
    box = slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                   Inches(12.20), Inches(0.40 * len(lines) + 0.3))
    set_lines(box, lines, size=size)
    return box


def _table(slide, rows, top, widths, height, size=12.5):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.60), Inches(top),
                                 Inches(12.20), Inches(height)).table
    for col, w in zip(tbl.columns, widths):
        col.width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            if not para.runs:
                continue
            para.runs[0].font.size = Pt(size)
            para.runs[0].font.bold = (r == 0)
            para.runs[0].font.color.rgb = NAVY
    return tbl


def slide_method(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "The Method That Came Out of Being Wrong", page)
    _body(s, 1.15, [
        "Section 7.10 was written three times and two claims withdrawn — each "
        "formed after the numbers were seen.",
        "Everything attempted afterwards fixed the test first. What that refused:",
    ])
    _table(s, [
        ("Fixed in advance", "What it refused"),
        ("7.15  the permutation null, built before reading the winner",
         "a +0.018 AUC gain over this work's headline readout"),
        ("7.16  the analysis, written while the arms were still folding",
         "a half-folded arm being scored at AUC 0.917"),
        ("7.17  the equivalence oracle, run before any timing",
         "three rewrites that looked correct and were slower or inexact"),
        ("7.19  the prediction, committed to source before a fold ran",
         "the option of calling the outcome expected after the fact"),
    ], 2.20, (6.4, 5.8), height=2.55)
    add_callout(s, 5.05,
                "The order matters more than the rigour. Each could have been done "
                "afterwards, produced the same arithmetic, and supported a weaker "
                "claim.", fill=GREEN)
    add_note(s, 6.05,
             "What pre-specification buys is not accuracy but the ability to be "
             "refused. In this work it refused an ensembling mechanism, a readout "
             "combination, a target-cropping strategy and a per-arm score — each "
             "of which looked correct when it was proposed.")
    return s


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src", default="final_dissertation_presentation.pptx")
    ap.add_argument("--out", dest="dst", default="final_dissertation_presentation.pptx")
    args = ap.parse_args()

    prs = Presentation(args.src)
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)
    if any("Came Out of Being Wrong" in t for t in titles):
        raise SystemExit("slide already present; nothing to do")
    try:
        insert_at = next(i for i, t in enumerate(titles) if t.startswith("Future Plan"))
    except StopIteration:
        raise SystemExit("no 'Future Plan' slide found; deck layout has changed")

    slide_method(prs, insert_at + 1)
    move_slide(prs, len(prs.slides) - 1, insert_at)
    print(f"  inserted at position {insert_at + 1}")

    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                    and sh.left and sh.left > Inches(12)):
                set_lines(sh, [str(i)], size=10, color=NAVY)

    prs.save(args.dst)
    print(f"wrote {args.dst}: {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
