"""Two places where the deck says less than the report now knows.

Slide 30 presents AUC 0.943 and concludes that the negative results are
properties of the regime rather than of the method. Section 7.18 has since
measured the same readout at 0.626 against binding that was actually measured,
and the counterweight slide sits five slides later -- so slide 30 read alone
gives the pre-7.18 story with nothing pointing forward.

The Summary ends at the screening tool and "extend the panel". The wet-lab check
and the settings decomposition, both of which change the practical claim, are
absent from it.

Neither is rewritten wholesale. Slide 30 keeps its result and gains a pointer;
the Summary gains two bullets and loses none.

Usage:
    python src/fix_presentation_consistency.py
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "src"))
from update_presentation import NAVY, set_lines  # noqa: E402


def find_slide(prs, needle):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and needle in sh.text_frame.text:
                return i, s, sh
    return None, None, None


def fix_caveat(prs):
    """Point slide 30's caveat at the section that qualifies it.

    The slide is already full to 6.90 in and the footer sits at 7.08, so the
    longer caveat does not fit where the short one did: written in place it ran
    to 7.33 and collided with the footer. The callout above it moves up and
    tightens to make room, and the caveat is kept to one line.
    """
    old = "Caveat: one model, 22 receptors, a single draw."
    i, s, sh = find_slide(prs, old)
    if sh is None:
        return False
    callout = next((x for x in s.shapes if x.has_text_frame
                    and "properties of the regime" in x.text_frame.text), None)
    if callout is not None:
        callout.top, callout.height = Inches(5.55), Inches(0.95)
    sh.top, sh.height = Inches(6.58), Inches(0.44)
    set_lines(sh, ["Caveat: one model, 22 receptors, one draw — and 0.943 is an "
                   "in-training figure; against measured binding it is 0.626 "
                   "(slide 35)."], size=12, color=NAVY)
    return True


def fix_summary(prs):
    """Add what the Summary predates, without disturbing what it already says."""
    i, s, sh = find_slide(prs, "Implemented and integrated memory")
    if sh is None:
        return False
    tf = sh.text_frame
    src = None
    for para in tf.paragraphs:
        if para.text.strip().startswith("•") and para.runs:
            src = para.runs[0].font
            break
    additions = [
        "•  Of the three inference settings raised together, sampling steps "
        "carry the effect; alignment depth and recycling carry none of it, and "
        "69% of the gain costs a third of the compute.",
        "•  Checked against binding measured in a wet lab on 1,320 designs no "
        "model had seen: 0.626, against 0.943 on the in-training panel. The "
        "readout still beats whole-chain confidence, which was the claim.",
    ]
    # insert before the closing "Next:" bullet so the slide still ends forward-looking
    paras = list(tf.paragraphs)
    nxt = next((p for p in paras if p.text.strip().startswith("•  Next:")), None)
    for text in additions:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = text
        if src is not None:
            r.font.size = src.size
            r.font.bold = src.bold
            if src.color and src.color.type is not None:
                r.font.color.rgb = src.color.rgb
    if nxt is not None:
        # move the Next: bullet to the end by rewriting it last
        body = nxt._p
        body.getparent().remove(body)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = nxt.text
        if src is not None:
            r.font.size = src.size
            r.font.bold = src.bold
            if src.color and src.color.type is not None:
                r.font.color.rgb = src.color.rgb
    return True


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src", default="final_dissertation_presentation.pptx")
    ap.add_argument("--out", dest="dst", default="final_dissertation_presentation.pptx")
    args = ap.parse_args()
    prs = Presentation(args.src)
    print(f"  caveat pointer : {'added' if fix_caveat(prs) else 'NOT FOUND'}")
    print(f"  summary bullets: {'added' if fix_summary(prs) else 'NOT FOUND'}")
    prs.save(args.dst)
    print(f"wrote {args.dst}: {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
