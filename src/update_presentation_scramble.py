"""Add Section 7.19 to the deck: the scramble control against measured binding.

The deck presents the scramble control as a requirement -- slide 31 has it built
into the screening tool, and the Summary lists it among what was delivered.
Section 7.19 puts a condition on it, tested against binding two CROs measured,
and a deck that states the rule without the condition now overstates it.

One slide, placed after the wet-lab slide it depends on. Its point is not that
the control failed but that its scope was measured, with the prediction
registered before the folds ran.

Usage:
    python src/update_presentation_scramble.py
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "src"))

from update_presentation import (  # noqa: E402
    AMBER, BLUE, GREEN, NAVY, add_callout, add_chrome, add_note,
    add_stat_card, blank_slide, move_slide, set_lines,
)


def _body(slide, top, lines, size=15):
    box = slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                   Inches(12.20), Inches(0.40 * len(lines) + 0.3))
    set_lines(box, lines, size=size)
    return box


def _table(slide, rows, top, widths, height=1.5, size=13):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.60), Inches(top),
                                 Inches(12.20), Inches(height)).table
    for col, w in zip(tbl.columns, widths):
        col.width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            if not para.runs:
                continue
            para.runs[0].font.size = Pt(size)
            para.runs[0].font.bold = (r == 0)
            para.runs[0].font.color.rgb = NAVY
    return tbl


def slide_scramble(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Measured — Where the Scramble Control Stops Working", page)
    _body(s, 1.15, [
        "38 of those designs folded as delivered and against permutations of "
        "themselves. The prediction below was",
        "written into the source and committed before any fold ran.",
    ])
    _table(s, [
        ("", "design", "its scrambles", "margin"),
        ("measured binders (20)", "72.99", "53.88", "+19.11"),
        ("measured non-binders (18)", "71.15", "53.22", "+17.93"),
    ], 2.15, (4.4, 2.6, 2.8, 2.4), height=1.35)
    add_stat_card(s, 0.60, 3.75, 3.85, "p = 0.751",
                  "the two margins are indistinguishable", AMBER)
    add_stat_card(s, 4.75, 3.75, 3.85, "0.672 → 0.581",
                  "AUC, before and after subtracting scrambles", AMBER)
    add_stat_card(s, 8.90, 3.75, 3.90, "−0.092",
                  "change in AUC, CI [−0.170, +0.008]", BLUE)
    add_callout(s, 5.50,
                "A permutation of a 15-residue peptide is still a plausible "
                "ligand. A permutation of a 100-residue protein is not a protein.",
                fill=GREEN)
    add_note(s, 6.05,
             "So the control is sound for the case this work concerns and unsound "
             "outside it. The recommendation gains a condition rather than losing "
             "its force: use the scramble control wherever a permutation of the "
             "candidate is still a candidate.")
    return s


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src", default="final_dissertation_presentation.pptx")
    ap.add_argument("--out", dest="dst", default="final_dissertation_presentation.pptx")
    args = ap.parse_args()

    prs = Presentation(args.src)
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)
    if any("Where the Scramble Control Stops Working" in t for t in titles):
        raise SystemExit("slide already present; nothing to do")
    # it depends on the wet-lab slide, so it goes immediately after it
    try:
        after = next(i for i, t in enumerate(titles)
                     if "Against Binding That Was Actually Measured" in t)
    except StopIteration:
        raise SystemExit("wet-lab slide not found; deck layout has changed")

    slide_scramble(prs, after + 2)
    move_slide(prs, len(prs.slides) - 1, after + 1)
    print(f"  inserted at position {after + 2}")

    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                    and sh.left and sh.left > Inches(12)):
                set_lines(sh, [str(i)], size=10, color=NAVY)

    prs.save(args.dst)
    print(f"wrote {args.dst}: {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
