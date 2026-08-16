"""Add the screening tool to the deck, before Future Plan.

The deck ends on measurement and then jumps to what is left to do. What is
missing is the thing the measurements were for: a tool a bench scientist can
run, with the project's three negatives built into its behaviour rather than
written in a caveat somewhere.

Two slides:

  * the tool itself -- what each finding turned into, and what it costs to run
  * the evidence it discriminates: the same three candidates against two
    targets, giving opposite answers, plus the target where it fails

Styling helpers come from update_presentation.py so the deck keeps one visual
language. Shape lookups are by text so this fails loudly rather than silently
rewriting the wrong box.

Usage:
    python src/update_presentation_tool.py \
        --in final_dissertation_presentation.pptx \
        --out final_dissertation_presentation.pptx
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "src"))

from update_presentation import (  # noqa: E402
    BLUE, GREEN, NAVY, TEAL, add_callout, add_chrome, add_note,
    add_stat_card, blank_slide, move_slide, set_lines,
)


def _body(slide, top, lines, size=15):
    box = slide.shapes.add_textbox(Inches(0.60), Inches(top),
                                   Inches(12.20), Inches(0.40 * len(lines) + 0.3))
    set_lines(box, lines, size=size)
    return box


def _table(slide, rows, top, widths, height=1.5, size=13):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.60), Inches(top),
                                 Inches(12.20), Inches(height)).table
    for col, w in zip(tbl.columns, widths):
        col.width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            if not para.runs:                    # an empty cell has nothing to style
                continue
            para.runs[0].font.size = Pt(size)
            para.runs[0].font.bold = (r == 0)
            para.runs[0].font.color.rgb = NAVY
    return tbl


def slide_tool(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "Delivered — the Findings as Behaviour, Not a Caveat", page)
    _body(s, 1.18, [
        "A local screening tool: paste a target and candidate peptides, get a ranking. "
        "Nothing is uploaded.",
    ])
    add_stat_card(s, 0.60, 1.90, 3.85, "every candidate",
                  "folded against permutations of itself", BLUE)
    add_stat_card(s, 4.75, 1.90, 3.85, "45 s + 9 s/fold",
                  "model construction is paid once per job", TEAL)
    add_stat_card(s, 8.90, 1.90, 3.90, "≈ 2 s",
                  "to re-screen — folds are cached", GREEN)
    items = [
        ("Section 7.4 — a score can be blind to sequence order",
         "so a candidate is reported as a t against its own permutations, never as a "
         "raw score; one that cannot beat its own scrambles is not a hit"),
        ("Section 7.13 — reduced settings suppress the effect 3–7×",
         "so quick mode is offered, and says in the interface what it costs: "
         "~14% of backbone bonds physically plausible at ten steps"),
        ("Section 7.5 — a single fold is not a measurement",
         "so replicates are a control, and the replicate spread is a column in the "
         "table rather than a footnote"),
    ]
    # the footer sits at 7.08 in and add_note is a fixed 0.95 in box, so the last
    # item has to finish by about 6.00 or the note runs into it
    top = 3.55
    for head, detail in items:
        box = s.shapes.add_textbox(Inches(0.60), Inches(top), Inches(12.20), Inches(0.80))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = head
        r0.font.size = Pt(14.5)
        r0.font.bold = True
        r0.font.color.rgb = NAVY
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = detail
        r1.font.size = Pt(12.5)
        r1.font.color.rgb = BLUE
        top += 0.82
    add_note(s, 6.02,
             "Candidates that cannot be scored honestly are refused, not scored badly. "
             "Building it caught two faults that would have misled a user: a "
             "two-permutation null SD reported z = +77, and a fixed cutoff on that ratio "
             "called a Gly–Ser linker a hit.")
    return s


def slide_two_targets(prs, page):
    s = blank_slide(prs)
    add_chrome(s, "The Same Three Candidates, Two Targets", page)
    _body(s, 1.18, [
        "Identical candidate list, identical settings, opposite answers — the readout "
        "is not simply rewarding peptide-shaped sequences.",
    ])
    rows = [
        ("candidate", "MDM2 — t", "MDM2 — p", "PDZ3 — t", "PDZ3 — p", "cognate of"),
        ("SQETFSDLWKLLPEN", "+3.88", "0.002", "+0.06", "0.477", "MDM2 (p53 helix)"),
        ("KQTSV", "+1.39", "0.101", "+7.72", "<0.001", "PSD-95 PDZ3"),
        ("PPPALPPKKR", "−0.00", "0.501", "+0.75", "0.240", "c-Crk SH3"),
    ]
    _table(s, rows, 2.00, (3.5, 1.6, 1.6, 1.6, 1.6, 2.3), height=1.7)
    _body(s, 3.85, [
        "Each cognate wins on its own target and loses on the other. KQTSV is also the "
        "shortest candidate at five",
        "residues, so the ranking is not being driven by length.",
    ])
    add_callout(s, 4.85,
                "Quick mode, two replicates, three scrambles — about three minutes a "
                "target on a laptop.", fill=GREEN)
    add_note(s, 6.02,
             "Shown with its failure: on c-Crk SH3 the cognate does not separate from "
             "its own scrambles (p = 0.40), and full sampling does not rescue it. A "
             "control that destroys only order has little to detect on a "
             "composition-driven interface.")
    return s


def correct_summary(prs):
    """Close the gap the new slides open in the Summary.

    Its last bullet still offers few-step distillation as future work, which the
    Future Plan slide has marked DONE since the DeCAF arm landed, and it lists
    the project's outputs without the one a reader can actually run. Both are
    replaced with a single pass so the deck does not summarise itself wrongly.
    """
    stale = "Next: GPU-scale replicate folding, then few-step sampler distillation."
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame or stale not in shape.text_frame.text:
                continue
            for para in shape.text_frame.paragraphs:
                if stale not in para.text:
                    continue
                for run in para.runs:                # keep the first run's styling
                    run.text = ""
                # the bullet glyph lives in the run text on this deck, not in the
                # paragraph format, so it has to be written back with the line
                para.runs[0].text = (
                    "•  Delivered as a local screening tool: the scramble control, the "
                    "cost of reduced settings and the replicate spread are built into "
                    "what it does rather than written beside it (slides 31–32).")
                extra = shape.text_frame.add_paragraph()
                r = extra.add_run()
                r.text = ("•  Next: extend the panel past 22 receptors, which is what "
                          "now limits the open questions.")
                src = para.runs[0].font
                r.font.size = src.size
                r.font.bold = src.bold
                if src.color and src.color.type is not None:
                    r.font.color.rgb = src.color.rgb
                return True
    return False


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="src", default="final_dissertation_presentation.pptx")
    ap.add_argument("--out", dest="dst", default="final_dissertation_presentation.pptx")
    args = ap.parse_args()

    prs = Presentation(args.src)
    n_before = len(prs.slides)
    print(f"opened {args.src}: {n_before} slides")

    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)
    if any("Findings as Behaviour" in t for t in titles):
        raise SystemExit("tool slides are already in this deck; nothing to do")
    try:
        insert_at = next(i for i, t in enumerate(titles) if t.startswith("Future Plan"))
    except StopIteration:
        raise SystemExit("no 'Future Plan' slide found; deck layout has changed")

    for k, build in enumerate([slide_tool, slide_two_targets]):
        build(prs, insert_at + k + 1)
        move_slide(prs, len(prs.slides) - 1, insert_at + k)
        print(f"  inserted at position {insert_at + k + 1}")

    print(f"  summary bullet {'corrected' if correct_summary(prs) else 'NOT FOUND'}")

    # renumber the trailing slides that shifted
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if (sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                    and sh.left and sh.left > Inches(12)):
                set_lines(sh, [str(i)], size=10, color=NAVY)

    prs.save(args.dst)
    print(f"wrote {args.dst}: {len(prs.slides)} slides (+{len(prs.slides) - n_before})")


if __name__ == "__main__":
    main()
